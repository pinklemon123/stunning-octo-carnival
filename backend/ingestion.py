import os
import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import io

def parse_file(file_storage, filename):
    """
    Dispatcher function to parse uploaded files based on extension.
    Returns the extracted text.
    """
    ext = os.path.splitext(filename)[1].lower()
    content = file_storage.read()
    
    if ext == '.txt' or ext == '.md':
        return parse_text(content)
    elif ext == '.pdf':
        return parse_pdf(content)
    elif ext == '.docx':
        return parse_docx(content)
    elif ext == '.pptx':
        return parse_pptx(content)
    elif ext == '.html' or ext == '.htm':
        return parse_html(content)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def parse_text(content):
    try:
        return content.decode('utf-8')
    except UnicodeDecodeError:
        try:
            return content.decode('gbk')
        except UnicodeDecodeError:
            return content.decode('latin-1')

def parse_pdf(content):
    try:
        reader = PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PDF: {e}")
        return ""

def parse_docx(content):
    try:
        doc = Document(io.BytesIO(content))
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing DOCX: {e}")
        return ""

def parse_pptx(content):
    try:
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        return ""

def parse_html(content):
    try:
        soup = BeautifulSoup(content, 'html.parser')
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator='\n')
        # Break into lines and remove leading/trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)
        return text
    except Exception as e:
        print(f"Error parsing HTML: {e}")
        return ""

def scrape_url(url):
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        return parse_html(response.content)
    except Exception as e:
        print(f"Error scraping URL {url}: {e}")
        raise e
